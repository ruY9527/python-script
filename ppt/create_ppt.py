import collections
import collections.abc # 显式导入以适配 Python 3.11+
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- 配置区 ---
OUTPUT_FILE = "妊娠合并前置胎盘与MG病例讨论_40页.pptx"
THEME_COLOR = RGBColor(0, 51, 102)  # 深蓝色医用主题
# BG_IMAGE = "background.jpg" # 如果你有背景图，取消注释并在 add_slide 中使用

def add_styled_slide(prs, title_str, points):
    """添加带样式的幻灯片"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题样式
    title = slide.shapes.title
    title.text = title_str
    title_tf = title.text_frame.paragraphs[0]
    title_tf.font.bold = True
    title_tf.font.size = Pt(32)
    title_tf.font.color.rgb = THEME_COLOR

    # 内容样式
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Pt(12)
        p.level = 0
        p.font.size = Pt(18)

def generate_medical_ppt():
    prs = Presentation()
    
    # 1-4: 封面与开场 [cite: 1, 2, 3, 4, 6]
    slide_0 = prs.slides.add_slide(prs.slide_layouts[0])
    slide_0.shapes.title.text = "一例妊娠合并完全性前置胎盘\n伴重症肌无力的疑难病例讨论 [cite: 1]"
    slide_0.shapes.placeholders[1].text = "主讲人：李雯婷 | 2025-10-26 [cite: 2, 3]"
    
    add_styled_slide(prs, "引言：高风险的碰撞 [cite: 4]", ["当产科最凶险的并发症遇上神经内科棘手的自身免疫病 [cite: 4]", "我们的应对策略与思考 [cite: 5]"])
    add_styled_slide(prs, "汇报目录 [cite: 6]", ["01 相关知识 [cite: 8]", "02 病例信息 [cite: 9]", "03 相关护理 [cite: 10]", "04 病例讨论 [cite: 11]"])
    add_styled_slide(prs, "MDT 多学科团队协作", ["产科、神经内科、麻醉科 [cite: 139, 142]", "风湿免疫科、神经外科、ICU [cite: 148, 153, 184]"])

    # 5-12: 第一部分 相关知识 [cite: 17, 32, 38, 44]
    add_styled_slide(prs, "重症肌无力 (MG) 定义 [cite: 17]", ["自身抗体介导、T 淋巴细胞辅助、补体参与 [cite: 17]", "骨骼肌波动性无力为主要表现 [cite: 17]"])
    add_styled_slide(prs, "MG 常见症状表现 [cite: 12]", ["乏力、复视 [cite: 12]", "眼睑下垂、呼吸困难、吞咽困难 [cite: 12]"])
    add_styled_slide(prs, "MG 的治疗矩阵 [cite: 18-24]", ["胆碱酯酶抑制剂、血浆置换 [cite: 18, 20]", "免疫抑制剂、静注免疫球蛋白 [cite: 21, 22]", "糖皮质激素、胸腺切除术 [cite: 23, 24]"])
    add_styled_slide(prs, "胎盘植入谱系疾病 (PAS) [cite: 32]", ["胎盘绒毛不同程度粘附或侵入子宫肌层 [cite: 32]", "严重母体并发症之一 [cite: 32]"])
    add_styled_slide(prs, "PAS 的病理分类 [cite: 34-37]", ["1. 胎盘粘连 [cite: 35]", "2. 胎盘植入 [cite: 36]", "3. 穿透性胎盘植入 [cite: 37]"])
    add_styled_slide(prs, "前置胎盘的临床意义 [cite: 38]", ["妊娠28周后位置低于胎先露部 [cite: 38]", "阴道出血主因，增加PAS发生风险 [cite: 38]"])
    add_styled_slide(prs, "MG 与妊娠的相互影响 [cite: 44, 45]", ["20-45岁女性高发，正值生育年龄 [cite: 46]", "妊娠/分娩可能诱发病情波动 [cite: 49]"])
    add_styled_slide(prs, "MG 妊娠转归结局 [cite: 48-52]", ["病情可能出现：稳定、改善或恶化 [cite: 50, 51, 52]", "需关注抗体类型与病情控制 [cite: 45]"])

    # 13-22: 第二部分 病例信息 [cite: 60, 69, 84, 138, 159]
    add_styled_slide(prs, "患者基本资料 [cite: 60-66]", ["姓名：韦某某 | 年龄：37岁 [cite: 62, 64]", "主诉：35+5周发现前置胎盘伴植入5月余 [cite: 68]"])
    add_styled_slide(prs, "入院主要诊断 [cite: 69-78]", ["1. 完全性前置胎盘伴胎盘植入 [cite: 70]", "2. 重症肌无力、系统性硬化症 [cite: 74, 75]", "3. 疤痕子宫、脐带绕颈 [cite: 72, 73]"])
    add_styled_slide(prs, "既往史：系统性硬化症 [cite: 84, 85]", ["病史4年，规律服用甲泼尼龙 4mg qd [cite: 85]", "辅助服用钙剂、维生素E/D [cite: 85]"])
    add_styled_slide(prs, "既往史：重症肌无力 [cite: 86]", ["病史4年，规律服用他克莫司 1mg bid [cite: 86]"])
    add_styled_slide(prs, "生育及手术史 [cite: 87-91]", ["异位妊娠切除术(2012)、剖宫产术(2014) [cite: 87, 88]", "人流术(2018)、宫腔镜手术(2022) [cite: 89, 90]"])
    add_styled_slide(prs, "入院体格检查 [cite: 97-102]", ["T:36.6℃, P:97次/分, BP:111/72mmHg [cite: 100]", "宫高33cm，胎心136次/分，无宫缩 [cite: 102]"])
    add_styled_slide(prs, "影像学检查 (MRI/超声) [cite: 106, 129, 132]", ["MRI: 晚孕、头位，考虑完全性前置胎盘并植入 [cite: 106]", "胎盘绒毛间隙狭窄，绒毛粘连 [cite: 131]"])
    add_styled_slide(prs, "实验室异常指标 [cite: 112-126]", ["血红蛋白: 101g/L (降低) [cite: 119]", "白蛋白: 24.7g/L (明显降低) [cite: 125, 126]"])
    add_styled_slide(prs, "MDT 意见：神内与麻醉 [cite: 138-142]", ["神内：关注麻醉镇静药，防范肌无力危象 [cite: 139]", "麻醉：备血、有创压监测、ICU备床 [cite: 142]"])
    add_styled_slide(prs, "主要用药方案 [cite: 154-172]", ["免疫调节：他克莫司、甲泼尼龙 [cite: 161, 169]", "预防支持：依诺肝素钠、钙片/维D [cite: 163, 171]"])

    # 23-30: 第三部分 相关护理 [cite: 173, 190, 218, 235]
    add_styled_slide(prs, "住院时间轴经过 [cite: 173-189]", ["10-13 入院 -> 10-15 MDT会诊 -> 10-16 手术 [cite: 178-185]", "10-21 好转出院 [cite: 189]"])
    add_styled_slide(prs, "入院风险评估 [cite: 190-207]", ["血栓评分: 4分 (高危) [cite: 201, 207]", "自理评估: 完全自理(100) [cite: 197, 203]"])
    add_styled_slide(prs, "术前护理诊断 (1/2) [cite: 218-225]", ["1. 潜在并发症：产前出血 [cite: 219]", "2. 躯体活动障碍/受伤风险 [cite: 221]"])
    add_styled_slide(prs, "术前护理诊断 (2/2) [cite: 218-225]", ["3. 潜在并发症：血栓 [cite: 223]", "4. 焦虑 (担心母婴预后) [cite: 225, 229]"])
    add_styled_slide(prs, "产前出血护理措施 [cite: 238-245]", ["绝对卧床、避免增加腹压、Q4h监测宫缩 [cite: 242, 243]", "应急演练：熟悉大出血抢救流程 [cite: 244]"])
    add_styled_slide(prs, "肌无力安全防护 [cite: 254-260]", ["定时评估呼吸肌/吞咽肌 [cite: 259]", "防跌倒标识，呼叫器触手可及 [cite: 258]"])
    add_styled_slide(prs, "血栓高危护理 [cite: 269-276]", ["指导踝泵运动、使用医用弹力袜 [cite: 274]", "皮下注射依诺肝素钠，每日测量腿围 [cite: 275, 276]"])
    add_styled_slide(prs, "心理与信息支持 [cite: 285-291]", ["通俗解释 MDT 优势，建立家庭支持联盟 [cite: 290, 291]"])

    # 31-36: 第四部分 手术与术后 [cite: 297, 308, 343, 372]
    add_styled_slide(prs, "手术详述 (10-16) [cite: 297, 299, 300]", ["腰硬联合麻醉下行子宫下段剖宫产术 [cite: 300, 301]", "术中胎盘植入达浆膜层，手取胎盘 [cite: 297]"])
    add_styled_slide(prs, "术中止血与转归 [cite: 297, 302]", ["出血量1000ml，输血2u红细胞+200ml血浆 [cite: 297, 302]", "放置球囊压迫止血，术后转ICU [cite: 297]"])
    add_styled_slide(prs, "术后转回病房评估 [cite: 308-332]", ["产后24h总出血 1125ml [cite: 309]", "血红蛋白降至 73g/L，自理能力中度依赖 [cite: 311, 328]"])
    add_styled_slide(prs, "术后护理挑战 (1/2) [cite: 343-352]", ["并发症预警：产后出血、肌无力危象 [cite: 344]", "感染风险：手术创伤及导管 [cite: 346]"])
    add_styled_slide(prs, "术后护理挑战 (2/2) [cite: 343-352]", ["母乳喂养困难、活动耐力下降 [cite: 348, 350]", "营养失调：低于机体需要量 [cite: 352]"])
    add_styled_slide(prs, "危象识别：护理重中之重 [cite: 372]", ["观察呼吸频率(快速浅呼吸)、锁骨上窝凹陷 [cite: 483]", "单次呼吸计数测试 (低于15提示差) [cite: 483]"])

    # 37-40: 第五部分 讨论与总结 [cite: 434, 469, 499]
    add_styled_slide(prs, "讨论：如何防范致命性产前出血？ [cite: 447, 458]", ["严密监测出血量、择期于孕36-37周手术 [cite: 452, 462]", "多团队协作(产科/麻醉/新生儿/神内) [cite: 463]"])
    add_styled_slide(prs, "讨论：MG 围术期管理核心 [cite: 469, 479]", ["使用 MG-ADL 量表量化评估症状 [cite: 470]", "避免应用禁忌药物：静脉注射镁剂 [cite: 485, 498]"])
    add_styled_slide(prs, "2025新版指南视角 [cite: 472, 473]", ["MSE (最小症状表达) 治疗目标：评分0-1分 [cite: 471]", "激素剂量维持在 5-10mg/天 达标状态 [cite: 497]"])
    add_styled_slide(prs, "病例小结：母婴平安 [cite: 499]", ["精心护理与 MDT 紧密配合是成功的关键 [cite: 495, 499]"])

    prs.save(OUTPUT_FILE)
    print(f"PPT 生成完成：{OUTPUT_FILE}")

if __name__ == "__main__":
    generate_medical_ppt()